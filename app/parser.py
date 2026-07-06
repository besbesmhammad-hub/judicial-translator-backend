from __future__ import annotations

import io
import re
import zipfile
from html import unescape

import fitz
from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader


def normalize(text: str) -> str:
    return (
        text.replace("\r\n", "\n")
        .replace("\x00", "")
        .replace("\u00a0", " ")
        .strip()
    )


def annotate_structure(text: str) -> str:
    lines = []
    for line in normalize(text).splitlines():
        trimmed = line.strip()
        if not trimmed:
            lines.append("")
        elif re.match(r"^(article|clause|section|chapitre)\s+[\w\dIVXLCDM.-]+", trimmed, re.I):
            lines.append(f"[ARTICLE] {trimmed}")
        elif re.match(r"^[A-ZÀ-Ÿ0-9 .,'()/-]{8,}$", trimmed) and len(trimmed) < 120:
            lines.append(f"[HEADING] {trimmed}")
        else:
            lines.append(line)
    return "\n".join(lines).strip()


def rtl_score(text: str) -> int:
    return len(re.findall(r"[\u0600-\u06FF]", text or ""))


def page_is_rtl(blocks: list[tuple]) -> bool:
    text = "\n".join(str(block[4]) for block in blocks if len(block) > 4)
    rtl = rtl_score(text)
    latin = len(re.findall(r"[A-Za-z]", text or ""))
    return rtl > 20 and rtl >= latin


def sorted_pdf_blocks(blocks: list[tuple]) -> list[tuple]:
    readable = [block for block in blocks if len(block) > 4 and str(block[4]).strip()]
    if page_is_rtl(readable):
        return sorted(readable, key=lambda block: (round(block[1], 1), -round(block[0], 1)))
    return sorted(readable, key=lambda block: (round(block[1], 1), round(block[0], 1)))


def parse_html(content: bytes | str) -> str:
    html = content.decode("utf-8", errors="ignore") if isinstance(content, bytes) else content
    soup = BeautifulSoup(html, "lxml")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()

    blocks: list[str] = []

    def table_to_markdown(table) -> str:
        rows = []
        for tr in table.find_all("tr"):
            cells = [cell.get_text(" ", strip=True).replace("|", "/") for cell in tr.find_all(["th", "td"])]
            if any(cells):
                rows.append(cells)
        if not rows:
            return ""
        width = max(len(row) for row in rows)
        rows = [row + [""] * (width - len(row)) for row in rows]
        return "\n".join(
            ["[TABLE]", "| " + " | ".join(rows[0]) + " |", "| " + " | ".join(["---"] * width) + " |"]
            + ["| " + " | ".join(row) + " |" for row in rows[1:]]
            + ["[/TABLE]"]
        )

    for element in soup.body.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "table"], recursive=True) if soup.body else []:
        name = element.name.lower()
        if element.find_parent("table") and name != "table":
            continue
        if name.startswith("h"):
            blocks.append(f"[HEADING {name[1]}] {element.get_text(' ', strip=True)}")
        elif name == "table":
            table = table_to_markdown(element)
            if table:
                blocks.append(table)
        elif name == "li":
            blocks.append(f"- {element.get_text(' ', strip=True)}")
        else:
            value = element.get_text(" ", strip=True)
            if value:
                blocks.append(value)

    text = "\n\n".join(blocks) or soup.get_text("\n", strip=True)
    return normalize(unescape(text))


def parse_docx(content: bytes) -> str:
    doc = Document(io.BytesIO(content))
    blocks: list[str] = []
    for paragraph in doc.paragraphs:
        value = paragraph.text.strip()
        if not value:
            continue
        style = (paragraph.style.name or "").lower() if paragraph.style else ""
        if "heading" in style or "titre" in style or "title" in style:
            blocks.append(f"[HEADING] {value}")
        elif re.match(r"^(article|clause|section|chapitre)\s+", value, re.I):
            blocks.append(f"[ARTICLE] {value}")
        else:
            blocks.append(value)

    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.replace("\n", " ").strip().replace("|", "/") for cell in row.cells]
            if any(cells):
                rows.append(cells)
        if rows:
            width = max(len(row) for row in rows)
            rows = [row + [""] * (width - len(row)) for row in rows]
            blocks.extend(
                ["[TABLE]", "| " + " | ".join(rows[0]) + " |", "| " + " | ".join(["---"] * width) + " |"]
                + ["| " + " | ".join(row) + " |" for row in rows[1:]]
                + ["[/TABLE]"]
            )
    return normalize("\n\n".join(blocks))


def parse_pptx(content: bytes) -> str:
    presentation = Presentation(io.BytesIO(content))
    blocks: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_blocks: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                value = shape.text.strip()
                if value:
                    slide_blocks.append(value)
            if getattr(shape, "has_table", False):
                rows = []
                for row in shape.table.rows:
                    cells = [cell.text.replace("\n", " ").strip().replace("|", "/") for cell in row.cells]
                    if any(cells):
                        rows.append(cells)
                if rows:
                    width = max(len(row) for row in rows)
                    rows = [row + [""] * (width - len(row)) for row in rows]
                    slide_blocks.extend(
                        ["[TABLE]", "| " + " | ".join(rows[0]) + " |", "| " + " | ".join(["---"] * width) + " |"]
                        + ["| " + " | ".join(row) + " |" for row in rows[1:]]
                        + ["[/TABLE]"]
                    )
        if slide_blocks:
            blocks.append(f"[SLIDE {slide_index}]\n" + "\n\n".join(slide_blocks))
    return normalize("\n\n".join(blocks))


def parse_xlsx(content: bytes) -> str:
    workbook = load_workbook(io.BytesIO(content), data_only=False, read_only=True)
    blocks: list[str] = []
    for sheet in workbook.worksheets:
        rows: list[list[str]] = []
        for row in sheet.iter_rows():
            cells = []
            for cell in row:
                value = cell.value
                cells.append("" if value is None else str(value).replace("|", "/").strip())
            if any(cells):
                rows.append(cells)
        if not rows:
            continue
        width = max(len(row) for row in rows)
        rows = [row + [""] * (width - len(row)) for row in rows]
        blocks.extend(
            [f"[SHEET] {sheet.title}", "[TABLE]", "| " + " | ".join(rows[0]) + " |", "| " + " | ".join(["---"] * width) + " |"]
            + ["| " + " | ".join(row) + " |" for row in rows[1:]]
            + ["[/TABLE]"]
        )
    workbook.close()
    return normalize("\n\n".join(blocks))


def parse_pdf(content: bytes) -> str:
    pages: list[str] = []
    try:
        pdf = fitz.open(stream=content, filetype="pdf")
        for index, page in enumerate(pdf, start=1):
            blocks = sorted_pdf_blocks(page.get_text("blocks"))
            text = "\n".join(block[4].strip() for block in blocks)
            pages.append(f"[PAGE {index}]\n{annotate_structure(text)}")
        parsed = "\n\n".join(pages).strip()
        if parsed:
            return parsed
    except Exception:
        pages.clear()

    reader = PdfReader(io.BytesIO(content))
    for index, page in enumerate(reader.pages, start=1):
        pages.append(f"[PAGE {index}]\n{annotate_structure(page.extract_text() or '')}")
    return normalize("\n\n".join(pages))


def detect_file_format(filename: str, content: bytes) -> str:
    lower = (filename or "").lower()
    header = content[:16]
    if header.startswith(b"%PDF"):
        return "pdf"
    if header.startswith(b"PK"):
        try:
            with zipfile.ZipFile(io.BytesIO(content)) as archive:
                names = set(archive.namelist())
            if "word/document.xml" in names:
                return "docx"
            if "ppt/presentation.xml" in names:
                return "pptx"
            if "xl/workbook.xml" in names:
                return "xlsx"
        except zipfile.BadZipFile:
            pass
    if lower.endswith(".docx"):
        return "docx"
    if lower.endswith(".pptx"):
        return "pptx"
    if lower.endswith((".xlsx", ".xlsm")):
        return "xlsx"
    if lower.endswith(".pdf"):
        return "pdf"
    if lower.endswith((".html", ".htm")) or content[:256].lstrip().lower().startswith((b"<!doctype html", b"<html")):
        return "html"
    return "txt"


def parse_document(filename: str, content: bytes) -> tuple[str, str]:
    file_format = detect_file_format(filename, content)
    if file_format == "docx":
        return parse_docx(content), "DOCX detected and parsed server-side with paragraphs, headings and tables."
    if file_format == "pptx":
        return parse_pptx(content), "PPTX detected and parsed server-side with slides, text boxes and tables."
    if file_format == "xlsx":
        return parse_xlsx(content), "XLSX detected and parsed server-side with sheets, rows and cells."
    if file_format == "pdf":
        return parse_pdf(content), "PDF detected and parsed server-side with page markers and layout-aware block order."
    if file_format == "html":
        return parse_html(content), "HTML detected and parsed server-side with headings, lists and tables."
    return annotate_structure(content.decode("utf-8", errors="ignore")), "Plain text parsed server-side with structural annotation."
