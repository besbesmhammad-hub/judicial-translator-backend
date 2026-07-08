from __future__ import annotations

import io
import re
import asyncio
import os
from dataclasses import dataclass
from typing import Awaitable, Callable

import fitz
import pytesseract
from PIL import Image
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from openpyxl import load_workbook
from openpyxl.styles import Alignment
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pytesseract import Output
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas

from .renderer import arabic_display, register_pdf_font


TranslateBatch = Callable[[list[str], str], Awaitable[list[str]]]


@dataclass
class TextItem:
    text: str
    setter: Callable[[str], None]


def comparable_token(value: str) -> str:
    return re.sub(r"[,\s\u00A0]", "", str(value or "").upper())


def protected_tokens(source: str) -> list[str]:
    patterns = [
        r"\b(?:DTU|NF|ISO|API|PDF|PPTX|DOCX|XLSX|HTML|TVA|HT|TTC|IRPP|IS)\b(?:\s*\d+(?:[.,-]\d+)*)?",
        r"\bEurocode\s*\d+(?:[.,-]\d+)*",
        r"\b(?:article|art\.|loi|décret|decret|ordonnance|arrêté|arrete|clause|section)\s*[A-Z]?\d+(?:[.-]\d+)*",
        r"\d{1,2}[\/.-]\d{1,2}[\/.-]\d{2,4}",
        r"\d[\d\s.,]*(?:\s?(?:€|\$|%|HT|TTC|TVA))+",
    ]
    seen: set[str] = set()
    out: list[str] = []
    for pattern in patterns:
        for match in re.finditer(pattern, source or "", flags=re.I):
            token = match.group(0).strip()
            key = comparable_token(token)
            if token and key not in seen:
                seen.add(key)
                out.append(token)
    return out


def numeric_tokens(value: str) -> list[str]:
    return re.findall(r"\d[\d\s.,]*(?:\s?(?:%|€|\$|HT|TTC))?", value or "", flags=re.I)


def preserve_numeric_tokens(source: str, translated: str) -> str:
    source_nums = numeric_tokens(source)
    if not source_nums:
        return translated
    index = 0

    def repl(match: re.Match[str]) -> str:
        nonlocal index
        replacement = source_nums[index] if index < len(source_nums) else match.group(0)
        index += 1
        return replacement

    return re.sub(r"\d[\d\s.,]*(?:\s?(?:%|€|\$|HT|TTC))?", repl, translated or "", flags=re.I)


def restore_missing_protected_tokens(source: str, translated: str) -> str:
    target_key = comparable_token(translated)
    missing = [token for token in protected_tokens(source) if comparable_token(token) not in target_key]
    if not missing:
        return translated
    suffix = "، ".join(missing)
    return f"{translated} ({suffix})" if translated else suffix


def repair_arabic_terms(source: str, translated: str) -> str:
    value = translated or ""
    replacements = [
        (r"الاكتشاف", "المعاينة"),
        (r"الاضطرابات", "العيوب"),
        (r"تآكل التعزيزات", "تآكل حديد التسليح"),
        (r"المواصفات الفنية", "دفتر الشروط الفنية"),
        (r"المؤسسة", "الهيكل الإنشائي"),
        (r"الدرجة", "البند"),
        (r"غير شاملة الضريبة", "دون احتساب الرسوم"),
        (r"إصلاح التشخيص", "التشخيص ← الإصلاح"),
        (r"التشخيص\s*→\s*الإصلاح", "التشخيص ← الإصلاح"),
    ]
    for pattern, replacement in replacements:
        value = re.sub(pattern, replacement, value)
    return value


# Matches JT_SEG markers including malformed variants (2-3 brackets, stray spaces).
SEGMENT_LOOSE_RE = re.compile(r"\[{2,3}\s*/?\s*JT_SEG_\d{4}\s*\]{2,3}")


def strip_segment_markers(value: str) -> str:
    """Final safety net: remove any residual JT_SEG markers before text reaches the document."""
    return SEGMENT_LOOSE_RE.sub("", value or "").strip()


def finalize_native_translation(source: str, translated: str) -> str:
    value = strip_segment_markers(translated) or source
    if re.search(r"[\u0600-\u06FF]", value):
        value = repair_arabic_terms(source, value)
    value = preserve_numeric_tokens(source, value)
    value = restore_missing_protected_tokens(source, value)
    return value.strip()


def apply_docx_rtl(paragraph) -> None:
    paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    ppr = paragraph._p.get_or_add_pPr()
    bidi = ppr.find(qn("w:bidi"))
    if bidi is None:
        bidi = OxmlElement("w:bidi")
        ppr.append(bidi)
    bidi.set(qn("w:val"), "1")
    for run in paragraph.runs:
        run.font.rtl = True


def should_translate(value: str | None) -> bool:
    text = (value or "").strip()
    if len(text) < 2:
        return False
    if text.startswith("="):
        return False
    if re.fullmatch(r"[\W\d_]+", text, flags=re.UNICODE):
        return False
    if re.fullmatch(r"https?://\S+|\S+@\S+\.\S+", text, flags=re.I):
        return False
    return True


def replace_paragraph_text(paragraph, translated: str) -> None:
    if not paragraph.runs:
        paragraph.add_run(translated)
    else:
        paragraph.runs[0].text = translated
        for run in paragraph.runs[1:]:
            run.text = ""
    if re.search(r"[\u0600-\u06FF]", translated or ""):
        apply_docx_rtl(paragraph)


def iter_docx_paragraphs(container):
    for paragraph in getattr(container, "paragraphs", []):
        yield paragraph
    for table in getattr(container, "tables", []):
        for row in table.rows:
            for cell in row.cells:
                yield from iter_docx_paragraphs(cell)


def collect_docx_items(document: Document) -> list[TextItem]:
    items: list[TextItem] = []
    containers = [document]
    for section in document.sections:
        containers.extend([section.header, section.footer])

    for container in containers:
        for paragraph in iter_docx_paragraphs(container):
            text = paragraph.text
            if should_translate(text):
                items.append(TextItem(text=text, setter=lambda value, p=paragraph: replace_paragraph_text(p, value)))
    return items


def replace_pptx_paragraph_text(paragraph, translated: str) -> None:
    text_frame = getattr(getattr(paragraph, "_parent", None), "text_frame", None)
    if text_frame is not None:
        text_frame.word_wrap = True
        try:
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception:
            pass
    if not paragraph.runs:
        paragraph.add_run().text = translated
    else:
        paragraph.runs[0].text = translated
        for run in paragraph.runs[1:]:
            run.text = ""
    if re.search(r"[\u0600-\u06FF]", translated or ""):
        paragraph.alignment = PP_ALIGN.RIGHT
        ppr = paragraph._p.get_or_add_pPr()
        ppr.set("rtl", "1")
        for run in paragraph.runs:
            run.font.name = "Arial"
            try:
                if run.font.size:
                    run.font.size = int(run.font.size * 0.9)
            except Exception:
                pass


def collect_pptx_shape_items(shape) -> list[TextItem]:
    items: list[TextItem] = []
    if getattr(shape, "shape_type", None) == 6:
        for subshape in shape.shapes:
            items.extend(collect_pptx_shape_items(subshape))
        return items

    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            text = "".join(run.text for run in paragraph.runs).strip() or paragraph.text.strip()
            if should_translate(text):
                items.append(TextItem(text=text, setter=lambda value, p=paragraph: replace_pptx_paragraph_text(p, value)))

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip() or paragraph.text.strip()
                    if should_translate(text):
                        items.append(TextItem(text=text, setter=lambda value, p=paragraph: replace_pptx_paragraph_text(p, value)))
    return items


def collect_pptx_items(presentation: Presentation) -> list[TextItem]:
    items: list[TextItem] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            items.extend(collect_pptx_shape_items(shape))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame
            for paragraph in notes.paragraphs:
                text = "".join(run.text for run in paragraph.runs).strip() or paragraph.text.strip()
                if should_translate(text):
                    items.append(TextItem(text=text, setter=lambda value, p=paragraph: replace_pptx_paragraph_text(p, value)))
    return items


def collect_xlsx_items(workbook) -> list[TextItem]:
    items: list[TextItem] = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if isinstance(cell.value, str) and should_translate(cell.value):
                    items.append(TextItem(text=cell.value, setter=lambda value, c=cell: setattr(c, "value", value)))
    return items


async def apply_translations(items: list[TextItem], translate_batch: TranslateBatch, context: str) -> int:
    if not items:
        return 0
    source_texts = [item.text for item in items]
    translated = await translate_batch(source_texts, context)
    if len(translated) != len(items):
        raise ValueError("Native document translation returned a mismatched segment count.")
    for item, value in zip(items, translated):
        item.setter(finalize_native_translation(item.text, value))
    return len(items)


async def translate_docx_native(content: bytes, translate_batch: TranslateBatch) -> tuple[bytes, str, str, int]:
    document = Document(io.BytesIO(content))
    changed = await apply_translations(collect_docx_items(document), translate_batch, "DOCX native in-place translation")
    stream = io.BytesIO()
    document.save(stream)
    return (
        stream.getvalue(),
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "docx",
        changed,
    )


async def translate_pptx_native(content: bytes, translate_batch: TranslateBatch) -> tuple[bytes, str, str, int]:
    presentation = Presentation(io.BytesIO(content))
    changed = await apply_translations(collect_pptx_items(presentation), translate_batch, "PPTX native in-place translation")
    stream = io.BytesIO()
    presentation.save(stream)
    return (
        stream.getvalue(),
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "pptx",
        changed,
    )


async def translate_xlsx_native(content: bytes, translate_batch: TranslateBatch) -> tuple[bytes, str, str, int]:
    workbook = load_workbook(io.BytesIO(content))
    changed = await apply_translations(collect_xlsx_items(workbook), translate_batch, "XLSX native in-place translation")
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if isinstance(cell.value, str) and re.search(r"[\u0600-\u06FF]", cell.value):
                    current = cell.alignment.copy()
                    cell.alignment = Alignment(
                        horizontal="right",
                        vertical=current.vertical,
                        wrap_text=True,
                        text_rotation=current.text_rotation,
                        shrink_to_fit=current.shrink_to_fit,
                        indent=current.indent,
                        readingOrder=2,
                    )
    stream = io.BytesIO()
    workbook.save(stream)
    return (
        stream.getvalue(),
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "xlsx",
        changed,
    )


def group_ocr_lines(data: dict, scale: float) -> list[dict]:
    grouped: dict[tuple[int, int], list[int]] = {}
    for index, text in enumerate(data.get("text", [])):
        if not should_translate(text):
            continue
        try:
            conf = float(data.get("conf", ["-1"])[index])
        except (TypeError, ValueError):
            conf = -1
        if conf >= 0 and conf < 35:
            continue
        key = (
            int(data.get("block_num", [0])[index]),
            int(data.get("par_num", [0])[index]),
        )
        grouped.setdefault(key, []).append(index)

    lines: list[dict] = []
    for indexes in grouped.values():
        words = [str(data["text"][index]).strip() for index in indexes if str(data["text"][index]).strip()]
        if not words:
            continue
        left = min(int(data["left"][index]) for index in indexes) / scale
        top = min(int(data["top"][index]) for index in indexes) / scale
        right = max(int(data["left"][index]) + int(data["width"][index]) for index in indexes) / scale
        bottom = max(int(data["top"][index]) + int(data["height"][index]) for index in indexes) / scale
        text = " ".join(words)
        if should_translate(text):
            lines.append({"text": text, "left": left, "top": top, "right": right, "bottom": bottom})
    return lines


def meaningful_block_text(text: str) -> str:
    return re.sub(r"[\W\d_]+", "", text or "", flags=re.UNICODE)


def extract_pdf_text_blocks(page) -> list[dict]:
    blocks: list[dict] = []
    try:
        page_dict = page.get_text("dict")
    except Exception:
        page_dict = {"blocks": []}

    for raw_block in page_dict.get("blocks", []):
        if raw_block.get("type") != 0:
            continue
        for line in raw_block.get("lines", []):
            spans = [span for span in line.get("spans", []) if str(span.get("text") or "").strip()]
            if not spans:
                continue
            text = re.sub(r"\s+", " ", " ".join(str(span.get("text") or "") for span in spans)).strip()
            if not should_translate(text):
                continue
            x0 = min(float(span["bbox"][0]) for span in spans)
            y0 = min(float(span["bbox"][1]) for span in spans)
            x1 = max(float(span["bbox"][2]) for span in spans)
            y1 = max(float(span["bbox"][3]) for span in spans)
            blocks.append({
                "text": text,
                "left": x0,
                "top": y0,
                "right": x1,
                "bottom": y1,
            })

    if not blocks:
        for raw in page.get_text("blocks"):
            if len(raw) < 5:
                continue
            text = re.sub(r"\s+", " ", str(raw[4] or "")).strip()
            if not should_translate(text):
                continue
            blocks.append({
                "text": text,
                "left": float(raw[0]),
                "top": float(raw[1]),
                "right": float(raw[2]),
                "bottom": float(raw[3]),
            })
    compact = "".join(meaningful_block_text(block["text"]) for block in blocks)
    return blocks if len(compact) >= 25 else []


def extract_pdf_page_blocks(page) -> list[dict]:
    blocks: list[dict] = []
    for raw in page.get_text("blocks"):
        if len(raw) < 5:
            continue
        text = re.sub(r"\s+", " ", str(raw[4] or "")).strip()
        if not should_translate(text):
            continue
        blocks.append({
            "text": text,
            "left": float(raw[0]),
            "top": float(raw[1]),
            "right": float(raw[2]),
            "bottom": float(raw[3]),
        })
    compact = "".join(meaningful_block_text(block["text"]) for block in blocks)
    return blocks if len(compact) >= 25 else []


def pdf_font_path() -> str | None:
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "C:/Windows/Fonts/arial.ttf",
        "C:/Windows/Fonts/Arial.ttf",
    ]
    return next((path for path in candidates if os.path.exists(path)), None)


def sampled_pdf_fill(page, rect: fitz.Rect) -> tuple[float, float, float]:
    try:
        pixmap = page.get_pixmap(matrix=fitz.Matrix(1, 1), alpha=False)
        image = Image.open(io.BytesIO(pixmap.tobytes("png"))).convert("RGB")
        width, height = image.size
        margin = 4
        points = [
            (rect.x0 - margin, rect.y0 - margin),
            (rect.x1 + margin, rect.y0 - margin),
            (rect.x0 - margin, rect.y1 + margin),
            (rect.x1 + margin, rect.y1 + margin),
            (rect.x0 - margin, (rect.y0 + rect.y1) / 2),
            (rect.x1 + margin, (rect.y0 + rect.y1) / 2),
            ((rect.x0 + rect.x1) / 2, rect.y0 - margin),
            ((rect.x0 + rect.x1) / 2, rect.y1 + margin),
        ]
        colors = []
        for x_raw, y_raw in points:
            x = max(0, min(width - 1, int(x_raw)))
            y = max(0, min(height - 1, int(y_raw)))
            colors.append(image.getpixel((x, y)))
        channels = list(zip(*colors))
        return tuple(sorted(channel)[len(channel) // 2] / 255 for channel in channels)
    except Exception:
        return (1, 1, 1)


def fitz_align_for_text(text: str) -> int:
    return fitz.TEXT_ALIGN_RIGHT if re.search(r"[\u0600-\u06FF]", text or "") else fitz.TEXT_ALIGN_LEFT


def fitz_text_for_pdf(text: str) -> str:
    return arabic_display(text) if re.search(r"[\u0600-\u06FF]", text or "") else text


def text_color_for_fill(fill: tuple[float, float, float]) -> tuple[float, float, float]:
    r, g, b = fill
    luminance = 0.2126 * r + 0.7152 * g + 0.0722 * b
    return (1, 1, 1) if luminance < 0.42 else (0.04, 0.04, 0.04)


def insert_replacement_text(page, rect: fitz.Rect, text: str, font_name: str, fill: tuple[float, float, float]) -> None:
    value = fitz_text_for_pdf(text)
    align = fitz_align_for_text(text)
    color = text_color_for_fill(fill)
    box_height = max(8, rect.height)
    font_size = min(13, max(6, box_height * 0.68))
    padded = fitz.Rect(rect.x0 + 1, rect.y0 + 1, rect.x1 - 1, rect.y1 - 1)
    for size in [font_size, 11, 10, 9, 8, 7, 6]:
        if size > font_size:
            continue
        leftover = page.insert_textbox(
            padded,
            value,
            fontname=font_name,
            fontsize=size,
            color=color,
            align=align,
        )
        if leftover >= -0.5:
            return


async def translate_pdf_text_native(content: bytes, translate_batch: TranslateBatch) -> tuple[bytes, str, str, int] | None:
    doc = fitz.open(stream=content, filetype="pdf")
    pages: list[dict] = []
    items: list[dict] = []
    for page_index, page in enumerate(doc):
        blocks = extract_pdf_text_blocks(page)
        page_items = []
        for block in blocks:
            rect = fitz.Rect(block["left"], block["top"], block["right"], block["bottom"])
            item = {
                "text": block["text"],
                "rect": rect,
                "fill": sampled_pdf_fill(page, rect),
            }
            page_items.append(item)
            items.append({"page": page_index, **item})
        pages.append({"items": page_items})

    if not items:
        doc.close()
        return None

    translated = await translate_batch([item["text"] for item in items], "PDF embedded-text in-place translation")
    if len(translated) != len(items):
        doc.close()
        raise ValueError("PDF text translation returned a mismatched segment count.")

    by_page: dict[int, list[dict]] = {}
    for item, value in zip(items, translated):
        item["translated"] = finalize_native_translation(item["text"], value)
        by_page.setdefault(item["page"], []).append(item)

    font_file = pdf_font_path()
    for page_index, page in enumerate(doc):
        page_items = by_page.get(page_index, [])
        if not page_items:
            continue
        for item in page_items:
            page.add_redact_annot(item["rect"], fill=item["fill"])
        page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)
        font_name = "helv"
        if font_file:
            try:
                page.insert_font(fontname="JudicialSans", fontfile=font_file)
                font_name = "JudicialSans"
            except Exception:
                font_name = "helv"
        for item in page_items:
            insert_replacement_text(page, item["rect"], item["translated"], font_name, item["fill"])

    output = doc.tobytes(garbage=4, deflate=True)
    doc.close()
    return output, "application/pdf", "pdf", len(items)


def collect_pdf_visual_items(content: bytes) -> tuple[list[dict], list[dict]]:
    scale = 1.8
    pages: list[dict] = []
    items: list[dict] = []
    with fitz.open(stream=content, filetype="pdf") as pdf:
        for page_index, page in enumerate(pdf):
            rect = page.rect
            pixmap = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
            image_bytes = pixmap.tobytes("png")
            lines = extract_pdf_text_blocks(page)
            if not lines:
                image = Image.open(io.BytesIO(image_bytes))
                data = pytesseract.image_to_data(
                    image,
                    lang="ara+fra+eng",
                    config="--psm 6 -c preserve_interword_spaces=1",
                    output_type=Output.DICT,
                )
                lines = group_ocr_lines(data, scale)
            page_record = {
                "width": float(rect.width),
                "height": float(rect.height),
                "image": image_bytes,
                "lines": lines,
            }
            pages.append(page_record)
            for line in lines:
                items.append({"page": page_index, "line": line, "text": line["text"]})
    return pages, items


def draw_wrapped_overlay(pdf, text: str, box: dict, page_height: float, font_name: str) -> None:
    x = box["left"]
    y_top = box["top"]
    width = max(30, box["right"] - box["left"])
    height = max(12, box["bottom"] - box["top"])
    rtl = bool(re.search(r"[\u0600-\u06FF]", text or ""))
    font_size = max(6, min(13, height * 0.72))
    max_chars = max(8, int(width / max(font_size * 0.42, 3)))
    chunks = [text[index:index + max_chars] for index in range(0, len(text), max_chars)] or [text]

    y = page_height - y_top - font_size
    for chunk in chunks[: max(1, int(height // max(font_size, 1)) + 1)]:
        pdf.setFont(font_name, font_size)
        if rtl:
            pdf.drawRightString(x + width, y, arabic_display(chunk))
        else:
            pdf.drawString(x, y, chunk)
        y -= font_size + 2


def render_pdf_visual_overlay(pages: list[dict]) -> bytes:
    stream = io.BytesIO()
    first = pages[0]
    pdf = canvas.Canvas(stream, pagesize=(first["width"], first["height"]))
    font_name = register_pdf_font()
    for page_index, page in enumerate(pages):
        if page_index:
            pdf.setPageSize((page["width"], page["height"]))
        pdf.drawImage(ImageReader(io.BytesIO(page["image"])), 0, 0, width=page["width"], height=page["height"])
        for line in page["lines"]:
            translated_line = (line.get("translated") or "").strip()
            if not translated_line:
                continue
            x = line["left"]
            y = page["height"] - line["bottom"]
            w = max(1, line["right"] - line["left"])
            h = max(1, line["bottom"] - line["top"])
            pdf.setFillColorRGB(1, 1, 1)
            pdf.rect(x - 1, y - 1, w + 2, h + 3, stroke=0, fill=1)
            pdf.setFillColorRGB(0.05, 0.05, 0.05)
            draw_wrapped_overlay(pdf, translated_line, line, page["height"], font_name)
        pdf.showPage()
    pdf.save()
    return stream.getvalue()


async def translate_pdf_visual_native(content: bytes, translate_batch: TranslateBatch) -> tuple[bytes, str, str, int]:
    text_native = await translate_pdf_text_native(content, translate_batch)
    if text_native is not None:
        return text_native

    pages, items = await asyncio.to_thread(collect_pdf_visual_items, content)
    if not items:
        raise ValueError("No OCR text blocks found for visual PDF translation.")
    translated = await translate_batch([item["text"] for item in items], "PDF visual in-place overlay translation")
    if len(translated) != len(items):
        raise ValueError("Visual PDF translation returned a mismatched segment count.")
    for item, value in zip(items, translated):
        item["line"]["translated"] = finalize_native_translation(item["text"], value)

    content_out = await asyncio.to_thread(render_pdf_visual_overlay, pages)
    return content_out, "application/pdf", "pdf", len(items)
